import os
from pptx import Presentation

# Input and output folders
input_dir = "../data/course_content"
output_dir = "../data/course_content_extracted/ppt"
os.makedirs(output_dir, exist_ok=True)

def extract_text_from_ppt(ppt_path):
    text = ""
    try:
        prs = Presentation(ppt_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        print(f"Error reading {ppt_path}: {e}")
    return text

def process_all_ppts():
    for root, dirs, files in os.walk(input_dir):
        for file in files:
            if file.endswith(".pptx"):
                ppt_path = os.path.join(root, file)
                rel_path = os.path.relpath(ppt_path, input_dir)
                output_path = os.path.join(output_dir, rel_path.replace('.pptx', '.txt'))

                os.makedirs(os.path.dirname(output_path), exist_ok=True)

                print(f"Extracting: {ppt_path}")
                text = extract_text_from_ppt(ppt_path)

                with open(output_path, "w", encoding="utf-8") as f:
                    f.write(text)

                print(f"✅ Saved: {output_path}")

if __name__ == "__main__":
    process_all_ppts()
